import PyPDF2
import docx
from pptx import Presentation

class FileProcessor:
    def extract_text(self, file_path: str) -> str:
        """
        Determines the file type by extension and extracts text accordingly.
        
        :param file_path: Path to the uploaded file
        :return: Extracted text from the file or an error message
        """
        file_extension = file_path.split('.')[-1].lower()
        
        try:
            if file_extension == "pdf":
                return self._extract_text_from_pdf(file_path)
            elif file_extension in ["doc", "docx"]:
                return self._extract_text_from_docx(file_path)
            elif file_extension in ["ppt", "pptx"]:
                return self._extract_text_from_pptx(file_path)
            else:
                return "Unsupported file type. Please upload a PDF, DOCX, or PPTX file."
        except Exception as e:
            return f"An error occurred while processing the file: {str(e)}"

    def _extract_text_from_pdf(self, file_path: str) -> str:
        """
        Extracts text from a PDF file.

        :param file_path: Path to the PDF file
        :return: Extracted text from PDF
        """
        text = ""
        with open(file_path, "rb") as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text
        return text

    def _extract_text_from_docx(self, file_path: str) -> str:
        """
        Extracts text from a DOCX file.

        :param file_path: Path to the DOCX file
        :return: Extracted text from DOCX
        """
        text = ""
        try:
            doc = docx.Document(file_path)
            text = "\n".join([para.text for para in doc.paragraphs])
        except Exception as e:
            text = f"Error reading DOCX file: {str(e)}"
        return text

    def _extract_text_from_pptx(self, file_path: str) -> str:
        """
        Extracts text from a PPTX file.

        :param file_path: Path to the PPTX file
        :return: Extracted text from PPTX
        """
        text = ""
        try:
            ppt = Presentation(file_path)
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        except Exception as e:
            text = f"Error reading PPTX file: {str(e)}"
        return text